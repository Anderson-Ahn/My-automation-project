import os
from pptx import Presentation
from openpyxl import Workbook

# 폴더 내 모든 PPT 파일 처리
ppt_folder = "D:\Python\proposal"  # PPT 파일이 저장된 폴더 경로
ppt_files = [f for f in os.listdir(ppt_folder) if f.endswith(".pptx")]

# Excel 파일 생성
wb = Workbook()
ws = wb.active
ws.title = "PPT Data"

# 헤더 추가
ws.append(["File Name", "Abstract Data", "Keyword Data"])

# 각 PPT 파일 처리
for ppt_file in ppt_files:
    ppt_path = os.path.join(ppt_folder, ppt_file)
    presentation = Presentation(ppt_path)

    # 4번째 슬라이드 가져오기
    if len(presentation.slides) < 4:
        print(f"'{ppt_file}'에는 4번째 슬라이드가 없습니다. 건너뜁니다.")
        continue
    slide = presentation.slides[3]  # 슬라이드는 0부터 시작

    # 첫 번째 표 가져오기
    table = None
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            break

    if table is None:
        print(f"'{ppt_file}'의 4번째 슬라이드에 표가 없습니다. 건너뜁니다.")
        continue

    # 1번째 행과 2번째 행 데이터 가져오기
    try:
        row1 = [cell.text for cell in table.rows[0].cells]  # 첫 번째 행
        row2_col2 = table.rows[1].cells[1].text  # 두 번째 행
    except IndexError:
        print(f"'{ppt_file}'의 표에 필요한 행이 부족합니다. 건너뜁니다.")
        continue

    # 데이터를 하나의 문자열로 합치기
    row1_data = ", ".join(row1)

    # Excel에 데이터 추가
    ws.append([ppt_file, row1_data, row2_col2])

# Excel 파일 저장
excel_path = "textmining.xlsx"  # 저장할 Excel 파일 경로
wb.save(excel_path)

print(f"모든 데이터가 '{excel_path}'에 저장되었습니다.")
